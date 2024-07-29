import os
import win32com.client
from pptx import Presentation

def extract_text_from_ppt(ppt_file):
    if ppt_file.endswith('.ppt'):
        convert_ppt_to_pptx(ppt_file, 'Presentations\\temp.pptx')
        ppt_file = 'Presentations\\temp.pptx'
    prs = Presentation(ppt_file)
    extracted_text = []

    for slide in prs.slides:
        text_slide = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                shape_text = text_slide.append(shape.text.splitlines())
        text_slide = [line for line in text_slide if line is not None]
        text_slide = [line for line in text_slide if line != []]
        extracted_text.append(text_slide)
    extracted_text = [line for line in extracted_text if line is not None]
    if 'temp.pptx' in os.listdir('Presentations'):
        os.remove('Presentations\\temp.pptx')
    return extracted_text

def convert_ppt_to_pptx(ppt_file, pptx_file):
    # Create a PowerPoint application object
    powerpoint = win32com.client.Dispatch("PowerPoint.Application")

    # Open the existing PPT file
    presentation = powerpoint.Presentations.Open(os.path.abspath(ppt_file))
    # Save the presentation as PPTX
    pptx_file = os.path.abspath(pptx_file)
    presentation.SaveAs(pptx_file, FileFormat=24)  # 24 represents PPTX format

    # Close the presentation and PowerPoint application
    presentation.Close()
    powerpoint.Quit()
